from typing import Dict, Any, List
from backend.agents.base import BaseAgent
from pptx import Presentation
from pptx.util import Inches
import os
import json

class PPTAgent(BaseAgent):
    def __init__(self, job_id: str):
        super().__init__("ppt", job_id)

    async def run(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        topic = input_data.get("instruction")
        self.log_activity("ppt_generation_started", {"topic": topic})

        # 1. Generate Outline & Content
        prompt = f"""
        Create a 5-slide presentation outline and content for: {topic}
        Return JSON format:
        {{
            "slides": [
                {{"title": "Slide 1 Title", "content": ["Bullet 1", "Bullet 2"]}},
                ...
            ]
        }}
        """
        response = await self.call_llm([{"role": "user", "content": prompt}], temperature=0.3)
        
        try:
            clean_response = response.replace("```json", "").replace("```", "").strip()
            data = json.loads(clean_response)
        except:
            return {"status": "failed", "error": "Failed to parse slide content"}

        # 2. Create PPTX
        prs = Presentation()
        
        for slide_data in data.get("slides", []):
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data.get("title", "Untitled")
            tf = content.text_frame
            
            for point in slide_data.get("content", []):
                p = tf.add_paragraph()
                p.text = point

        # 3. Save Artifact
        filename = f"presentation_{self.job_id}.pptx"
        filepath = os.path.join("artifacts", filename) # Ensure artifacts dir exists
        os.makedirs("artifacts", exist_ok=True)
        prs.save(filepath)
        
        self.log_activity("ppt_created", {"path": filepath})
        return {"status": "success", "path": filepath}
